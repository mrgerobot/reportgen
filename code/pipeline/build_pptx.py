from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os
from multiprocessing import Pool, cpu_count
import json
import matplotlib
import matplotlib.pyplot as plt
from io import BytesIO
import numpy as np
from multiprocessing import Semaphore
from pathlib import Path
matplotlib.use("Agg")
import sys
import logging
import time
import shutil
from pptx import Presentation
from pathlib import Path
import os
import logging

log = logging.getLogger(__name__)

def determine_template(student: dict, assets_dir: Path) -> Path:
    rol = student.get("Rol")

    if rol in ("counseling", "compass-directo"):
        return assets_dir / "Template_Autoconocimiento 2.0 (counseling).pptx"
    else:
        return assets_dir / "Template_Autoconocimiento 2.0 (completo).pptx"

def map_placeholders(student: dict) -> dict:
    mapped = {}

    for key, value in student.items():
        placeholder = f"<<{key}>>"
        mapped[placeholder] = value

    return mapped

def generate_report(
    student: dict,
    pie_chart_path: Path,
    template_path: Path,
    output_pptx_path: Path,
):

    log.info("Generating PPTX")
    prs = Presentation(template_path)

    placeholders = map_placeholders(student)

    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue

            text = shape.text

            for key, val in placeholders.items():

                # IMAGE PLACEHOLDERS
                if isinstance(val, dict) and "image" in val:
                    if key in text:
                        image_path = Path(val["image"])
                        if not image_path.exists():
                            raise FileNotFoundError(image_path)

                        left, top, width, height = (
                            shape.left, shape.top, shape.width, shape.height
                        )

                        shape.text = ""
                        slide.shapes.add_picture(
                            str(image_path),
                            left,
                            top,
                            width=width,
                            height=height
                        )

                # STRING PLACEHOLDERS
                elif isinstance(val, str):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if key in run.text:
                                if key == "<<URL GRAFICO HOLLAND>>":
                                    run.text = ""
                                    slide.shapes.add_picture(
                                        str(pie_chart_path),
                                        shape.left,
                                        shape.top,
                                        shape.width,
                                        shape.height
                                    )
                                else:
                                    run.text = run.text.replace(key, val)

                # NUMERIC PLACEHOLDERS
                elif isinstance(val, (int, float)):
                    formatted = f"{round(val * 100, 1)}%"
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if key in run.text:
                                run.text = run.text.replace(key, formatted)

    # Save PPTX
    output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx_path)

    log.info(f"PPTX written → {output_pptx_path}")
    return output_pptx_path


