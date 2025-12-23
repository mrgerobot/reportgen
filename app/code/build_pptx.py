from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os
from multiprocessing import Pool, cpu_count
import json
import matplotlib
import matplotlib.pyplot as plt
from io import BytesIO
import numpy as np
from multiprocessing import Semaphore
from pathlib import Path
matplotlib.use("Agg")
import sys
import logging
import time
import shutil

# file imports
from chartgen import *
from data_processing import *

logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)s] %(message)s"
)

log = logging.getLogger(__name__)

BASE_DIR = Path(__file__).resolve().parent.parent

# path a las plantillas
AC = BASE_DIR / "plantillas" / "AC"
COUNSELING_EST = AC / "Template_Autoconocimiento 2.0 (counseling).pptx"
COUNSELING_PADRES = AC / "Template_Autoconocimiento 2.0 (Padres - Counseling).pptx"
GS_ACT_EST = AC / "Template_Autoconocimiento 2.0 (completo).pptx"
GS_ACT_PADRES = AC / "Template_Autoconocimiento 2.0 (PADRES).pptx"


def find_libreoffice() -> str:
    """
    Return the path to LibreOffice (soffice/libreoffice) in a cross-platform way.
    Raises a clear error if not found.
    """
    # Common executable names, in priority order
    candidates = [
        "soffice",          # Linux, macOS, Docker
        "libreoffice",      # Some Linux distros
        "soffice.exe",      # Windows (PATH)
    ]

    for cmd in candidates:
        path = shutil.which(cmd)
        if path:
            return path

    raise FileNotFoundError(
        "LibreOffice not found.\n"
        "Make sure LibreOffice is installed and available in PATH.\n"
        "On Linux/Docker: apt install libreoffice\n"
        "On macOS: brew install libreoffice or install the app\n"
        "On Windows: enable 'soffice.exe' in PATH"
    )

def expected_pdf_name(student):
    # Adjust this to your actual JSON field
    name = student.get("Nombre y Apellido") or student.get("nombre") or "output"
    return f"{name}_Reporte_CCR.pdf"

def determine_template(student):
    rol = student["Rol"]
    if rol == "counseling" or rol == "compass-directo":
        template_estudiante = COUNSELING_EST
        template_padres = COUNSELING_PADRES
    else:
        template_estudiante = GS_ACT_EST
        template_padres = GS_ACT_PADRES
    
    return template_estudiante, template_padres

def generate_report(student, pie_chart, template_path, placeholders):
    log.info("=== generate_report START ===")
    log.debug(f"Template path: {template_path}")
    log.debug(f"Pie chart input: {pie_chart} (type={type(pie_chart)})")

    prs = Presentation(template_path)
    log.info(f"Loaded presentation with {len(prs.slides)} slides")

    for slide_idx, slide in enumerate(prs.slides):
        log.debug(f"-- Slide {slide_idx}")

        for shape_idx, shape in enumerate(list(slide.shapes)):
            log.debug(f"Shape {shape_idx}: has_text_frame={shape.has_text_frame}")

            if not shape.has_text_frame:
                continue

            text = shape.text
            log.debug(f"Shape {shape_idx} initial text: {repr(text)}")

            for key, val in placeholders.items():
                log.debug(f"Processing placeholder {key} (type={type(val)})")

                # IMAGE PLACEHOLDERS
                if isinstance(val, dict) and "image" in val:
                    if key in text:
                        log.info(f"Inserting image for key {key}")
                        image_path = val["image"]
                        log.debug(f"Image path: {image_path}")

                        if not os.path.exists(image_path):
                            raise FileNotFoundError(f"Image not found: {image_path}")

                        left, top, width, height = (
                            shape.left, shape.top, shape.width, shape.height
                        )
                        shape.text = ""

                        slide.shapes.add_picture(
                            image_path,
                            left,
                            top,
                            width=width,
                            height=height
                        )

                # STRING PLACEHOLDERS
                elif isinstance(val, str):
                    for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        for r_idx, run in enumerate(paragraph.runs):
                            if key in run.text:
                                log.info(
                                    f"Replacing text on slide {slide_idx}, "
                                    f"shape {shape_idx}, run {r_idx}"
                                )
                                log.debug(f"Run text BEFORE: {repr(run.text)}")

                                if key == "<<URL GRAFICO HOLLAND>>":
                                    log.info("Inserting pie chart image")

                                    if not os.path.exists(pie_chart):
                                        raise FileNotFoundError(
                                            f"Pie chart not found: {pie_chart}"
                                        )

                                    run.text = ""
                                    slide.shapes.add_picture(
                                        pie_chart,
                                        shape.left,
                                        shape.top,
                                        shape.width,
                                        shape.height
                                    )
                                else:
                                    run.text = run.text.replace(key, val)

                                log.debug(f"Run text AFTER: {repr(run.text)}")

                # INTEGER PLACEHOLDERS
                elif isinstance(val, int):
                    new_val = "%" + str(round(val * 100, 1))
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if key in run.text:
                                log.info(f"Replacing int placeholder {key}")
                                run.text = run.text.replace(key, new_val)

                # FLOAT PLACEHOLDERS
                elif isinstance(val, float):
                    new_val = "%" + str(round(val * 100, 1))
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if key in run.text:
                                log.info(f"Replacing float placeholder {key}")
                                run.text = run.text.replace(key, new_val)

    log.info("Finished modifying slides")

    # ---------- SAVE FINAL PPTX ----------
    OUTPUT_DIR = BASE_DIR / "outputs" / "pptx"
    OUTPUT_DIR.mkdir(exist_ok=True)
    log.info(f"El path al output es {OUTPUT_DIR}")
    student_id = student["Nombre y Apellido"]
    final_pptx_path = OUTPUT_DIR / f"{student_id}.pptx"
    tmp_pptx_path = final_pptx_path.with_suffix(".tmp")

    log.info(f"Saving PPTX to temp path: {tmp_pptx_path}")

    prs.save(tmp_pptx_path)

    # Atomic replace (safe if re-running or crashing)
    tmp_pptx_path.replace(final_pptx_path)

    log.info(f"PPTX generated successfully: {final_pptx_path}")

    return final_pptx_path

def map_placeholders(student):
    mapped = {}
    for key, value in student.items():
        placeholder = f"<<{key}>>"

        if key.lower() == "img":
            # Special case: image placeholder becomes a dict
            mapped[placeholder] = {"image": value}
        else:
            mapped[placeholder] = value

    return mapped

def process_student(student):
    try:
    
        mapped = map_placeholders(student)
        template_estudiante, _ = determine_template(student)
        chartpath = CHARTS_DIR / f"{student['Nombre y Apellido']}.png"
        if chartpath.exists():
            pie_chart = str(chartpath)
        else:
            pie_chart = generate_graph(student, save_debug=True)
        pdf_path = generate_report(student, pie_chart, template_path=template_estudiante, placeholders=mapped)
        return {"student": student, "status": "success", "file": pdf_path}

    except Exception as e:
        return {"student": student["Nombre y Apellido"], "status": "error", "error": str(e)}


data = BASE_DIR / "data" / "data_url_actualizadas.json"

if __name__ == "__main__":

    with open(data, "r", encoding="utf-8") as f:
        students = json.load(f)

    workers = min(15, max(1, cpu_count() - 1))
    with Pool(processes=workers) as pool:
        results = pool.map(process_student, students)
    print("\nAll pending reports generated.")

